"""PL Church Sermon PPT Converter  (Before → After)

Usage:
    python convert_pptx.py 에베소서 제4강(before).pptx
    python convert_pptx.py input.pptx output.pptx --png
"""

import sys
import io
import re
import math
import os
import argparse

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# ── 색상 ──────────────────────────────────────────────────────────────────────
BG_COLOR   = RGBColor(0x1B, 0x20, 0x45)
HEADER_BG  = RGBColor(0x23, 0x2A, 0x56)
GOLD       = RGBColor(0xF5, 0xC5, 0x18)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

# ── 슬라이드 크기 (표준 16:9, 1920×1080 기준) ────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

# ── 폰트 ─────────────────────────────────────────────────────────────────────
FONT_CLASSIC_BOLD = "국립박물관문화재단클래식 Bold"
FONT_CLASSIC_MED  = "국립박물관문화재단클래식 Medium"
FONT_NOTO         = "Noto Sans KR"

MAX_BULLETS = 3   # 슬라이드당 최대 불렛 수


# ── Before PPTX 파서 ──────────────────────────────────────────────────────────

def _para_text(para) -> str:
    return "".join(run.text for run in para.runs).strip()

def _tf_by_name(slide, name):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == name:
            return shape.text_frame
    return None

def _split_series_title(full_text: str):
    m = re.search(r"(.*?제\s*\d+\s*강)", full_text)
    if m:
        return m.group(1).strip(), full_text[m.end():].strip()
    return "", full_text.strip()

def parse_title_slide(slide):
    title_tf = _tf_by_name(slide, "제목 1")
    sub_tf   = _tf_by_name(slide, "부제목 2")

    full_text = ""
    if title_tf:
        for para in title_tf.paragraphs:
            t = _para_text(para)
            if t:
                full_text += t

    series, lecture_title = _split_series_title(full_text)

    sub_paras = []
    if sub_tf:
        for para in sub_tf.paragraphs:
            t = _para_text(para)
            if t:
                sub_paras.append(t)

    pastor = "Rev. Dr. Paul Junghoon Lee"

    if not lecture_title and sub_paras:
        non_pastor   = [t for t in sub_paras if "Rev" not in t and "Dr." not in t]
        pastor_lines = [t for t in sub_paras if "Rev" in t or "Dr." in t]
        lecture_title = non_pastor[0] if non_pastor else ""
        if pastor_lines:
            pastor = pastor_lines[0]
    else:
        for t in sub_paras:
            if t:
                pastor = t
                break

    return series, lecture_title, pastor

import io as _io

_NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

def _has_highlight(run) -> bool:
    """run에 <a:highlight> 요소가 있으면 True"""
    rPr = run._r.find(f'{{{_NS_A}}}rPr')
    if rPr is None:
        return False
    return rPr.find(f'{{{_NS_A}}}highlight') is not None

def _para_runs(para) -> list:
    """단락을 [{text, gold}] run 목록으로 반환.
    원본에 <a:highlight>가 있는 run만 gold=True; 없으면 모두 white.
    """
    raw_runs = [r for r in para.runs if r.text]
    if not raw_runs:
        return []

    any_hl = any(_has_highlight(r) for r in raw_runs)

    if not any_hl:
        full = _para_text(para)
        if not full:
            return []
        return [{'text': full, 'gold': False}]

    # 인접한 동일 gold 값 run 병합
    merged = []
    for run in raw_runs:
        gold = _has_highlight(run)
        if merged and merged[-1]['gold'] == gold:
            merged[-1]['text'] += run.text
        else:
            merged.append({'text': run.text, 'gold': gold})
    return [r for r in merged if r['text']]

def _bullet_full_text(bullet: list) -> str:
    return ''.join(r['text'] for r in bullet)

def _extract_images(slide, scale_x=1.0, scale_y=1.0) -> list:
    """슬라이드의 모든 이미지를 위치·크기 정보와 함께 추출한다.
    scale_x/y: 소스 슬라이드 → After 슬라이드 좌표 변환 비율.
    """
    images = []
    for shape in slide.shapes:
        blips = shape._element.findall(f'.//{{{_NS_A}}}blip')
        for blip in blips:
            rId = blip.get(f'{{{_NS_R}}}embed')
            if rId and rId in slide.part.rels:
                try:
                    blob = slide.part.rels[rId].target_part.blob
                    images.append({
                        'bytes':  blob,
                        'left':   int(shape.left   * scale_x),
                        'top':    int(shape.top    * scale_y),
                        'width':  int(shape.width  * scale_x),
                        'height': int(shape.height * scale_y),
                    })
                except Exception:
                    pass
    return images

def parse_content_slide(slide, scale_x=1.0, scale_y=1.0):
    title_tf = _tf_by_name(slide, "제목 1")

    # "내용 개체 틀 2", "내용 개체 틀 4" 등 번호 무관하게 탐색
    content_tf = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name.startswith("내용 개체 틀"):
            content_tf = shape.text_frame
            break

    heading_parts = []
    if title_tf:
        for para in title_tf.paragraphs:
            t = _para_text(para)
            if t:
                heading_parts.append(t)
    heading = " ".join(heading_parts)

    bullets = []  # List[List[{text, gold}]]
    content_paras = []
    if content_tf:
        content_paras = [p for p in content_tf.paragraphs if _para_text(p)]

    # 제목 1 shape이 없으면 내용 첫 단락을 heading으로 사용
    if not heading and content_paras:
        heading = _para_text(content_paras[0])
        content_paras = content_paras[1:]

    for para in content_paras:
        runs = _para_runs(para)
        if runs:
            bullets.append(runs)

    images = _extract_images(slide, scale_x, scale_y)
    return heading, bullets, images

def parse_before(path: str) -> dict:
    prs = Presentation(path)
    all_slides = list(prs.slides)

    # 소스 슬라이드 크기가 After 기준(SLIDE_W×SLIDE_H)과 다를 경우 좌표 스케일링
    scale_x = int(SLIDE_W) / int(prs.slide_width)
    scale_y = int(SLIDE_H) / int(prs.slide_height)

    series, lecture_title, pastor = parse_title_slide(all_slides[0])

    content_slides = []
    for slide in all_slides[1:]:
        heading, bullets, images = parse_content_slide(slide, scale_x, scale_y)
        # 텍스트가 전혀 없는 슬라이드는 스킵 (이미 변환된 파일 등 이미지 전용 슬라이드)
        if heading or bullets:
            content_slides.append({"heading": heading, "bullets": bullets, "images": images})

    return {
        "title": {"series": series, "lecture_title": lecture_title, "pastor": pastor},
        "content_slides": content_slides,
    }




# ── After PPTX 빌더 ───────────────────────────────────────────────────────────

def _solid_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _add_rect(slide, left, top, width, height, fill_color: RGBColor):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def _add_textbox(slide, left, top, width, height,
                 text, font_name, size_pt,
                 bold=None, italic=None, color=WHITE,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font_name
    run.font.size  = Pt(size_pt)
    if bold    is not None: run.font.bold   = bold
    if italic  is not None: run.font.italic = italic
    run.font.color.rgb = color
    return box

def build_title_slide(prs, series: str, lecture_title: str, pastor: str,
                      cover_image: str = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if cover_image and os.path.exists(cover_image):
        # 표지.png 를 슬라이드 전체 배경으로 (성경책·별 장식 모두 포함됨)
        slide.shapes.add_picture(cover_image, 0, 0, SLIDE_W, SLIDE_H)
    else:
        # 표지 이미지 없으면 단색 배경 + 별 텍스트 대체
        _solid_bg(slide, BG_COLOR)
        _add_textbox(
            slide,
            SLIDE_W - Inches(1.5), Inches(0.3),
            Inches(1.2), Inches(0.9),
            "✦", FONT_CLASSIC_BOLD, 48,
            bold=False, color=GOLD, align=PP_ALIGN.CENTER,
        )

    # ── 시리즈명 (예: 에베소서 제4강)
    _add_textbox(
        slide,
        Inches(1.5), Inches(1.9),
        SLIDE_W - Inches(3), Inches(0.6),
        series, FONT_CLASSIC_MED, 22,
        color=WHITE, align=PP_ALIGN.CENTER,
    )

    # ── 강의 제목 (중앙, 크게)
    _add_textbox(
        slide,
        Inches(1.5), Inches(2.55),
        SLIDE_W - Inches(3), Inches(2.2),
        lecture_title, FONT_CLASSIC_BOLD, 64,
        bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )

    # ── 작은 별 + 목사님 이름
    pastor_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(4.85),
        SLIDE_W - Inches(3), Inches(0.9),
    )
    tf = pastor_box.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    star_run = p.add_run()
    star_run.text = "✦  "
    star_run.font.name  = FONT_CLASSIC_BOLD
    star_run.font.size  = Pt(14)
    star_run.font.color.rgb = GOLD

    name_run = p.add_run()
    name_run.text = pastor
    name_run.font.name   = FONT_CLASSIC_MED
    name_run.font.size   = Pt(24)
    name_run.font.italic = True
    name_run.font.color.rgb = WHITE

    return slide

def build_content_slide(prs, heading: str, bullets: list, images=None, is_continued=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, BG_COLOR)

    HEADER_H = Inches(1.10)

    # ── 전체 배경 rect (일부 뷰어에서 background fill 미렌더링 대비)
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_COLOR)

    # ── 헤더 배경 박스
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, HEADER_H, HEADER_BG)

    # ── 헤더 텍스트 (계속 슬라이드도 동일 제목 사용)
    head_label = heading
    _add_textbox(
        slide,
        Inches(0.45), Inches(0.10),
        SLIDE_W - Inches(0.9), HEADER_H - Inches(0.08),
        head_label, FONT_CLASSIC_BOLD, 28,
        bold=True, color=WHITE, align=PP_ALIGN.LEFT,
    )

    # ── 본문 불렛 영역
    body_top = HEADER_H + Inches(0.15)
    body_h   = SLIDE_H  - body_top - Inches(0.15)
    body_box = slide.shapes.add_textbox(
        Inches(0.45), body_top,
        SLIDE_W - Inches(0.8), body_h,
    )
    tf = body_box.text_frame
    tf.word_wrap = True

    # hanging indent: 두 번째 줄이 "-" 뒤 텍스트와 정렬되도록
    _MARL   =  int(Inches(0.60))
    _INDENT = -int(Inches(0.60))

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment    = PP_ALIGN.LEFT
        p.space_before = Pt(10) if i == 0 else Pt(28)
        p.space_after  = Pt(0)

        pPr = p._p.get_or_add_pPr()
        pPr.set("marL",   str(_MARL))
        pPr.set("indent", str(_INDENT))

        # run별 색상 처리
        for j, rd in enumerate(bullet):
            run = p.add_run()
            run.text = ("-  " + rd['text']) if j == 0 else rd['text']
            run.font.name  = FONT_NOTO
            run.font.size  = Pt(32)
            run.font.bold  = True
            run.font.color.rgb = GOLD if rd['gold'] else WHITE

    # ── 이미지 삽입 (헤더 아래로 위치 보정)
    if images:
        HEADER_H_EMU = int(HEADER_H)
        for img in images:
            top    = img['top']
            height = img['height']
            # 이미지가 헤더와 겹치면 상단을 헤더 아래로 내림
            if top < HEADER_H_EMU:
                shift  = HEADER_H_EMU - top
                top    = HEADER_H_EMU
                height = max(height - shift, Emu(914400))  # 최소 1인치
            slide.shapes.add_picture(
                _io.BytesIO(img['bytes']),
                img['left'], top, img['width'], height,
            )

    return slide


_CHARS_PER_LINE  = 22   # 32pt에서 슬라이드 너비에 들어가는 한글 글자 수 추정
_MAX_LINES_SLIDE = 7    # 슬라이드 본문에 들어갈 최대 줄 수

def _est_lines(bullet: list) -> int:
    return max(1, math.ceil(len(_bullet_full_text(bullet)) / _CHARS_PER_LINE))

def _split_bullets(bullets: list, max_per_slide: int) -> list:
    n = len(bullets)
    if n == 0:
        return [bullets]
    total_lines = sum(_est_lines(b) for b in bullets)
    # 불렛 2개 이하는 길어도 분리 안 함
    if n <= 2:
        return [bullets]
    # 불렛 수가 허용치 초과: 줄 수도 넘을 때만 동등 분할
    # (줄 수 허용 범위면 불렛이 많아도 한 슬라이드에 배치)
    if n > max_per_slide:
        if total_lines <= _MAX_LINES_SLIDE:
            return [bullets]
        num_chunks = (n + max_per_slide - 1) // max_per_slide
        base_size  = n // num_chunks
        remainder  = n % num_chunks
        chunks, idx = [], 0
        for i in range(num_chunks):
            size = base_size + (1 if i < remainder else 0)
            chunks.append(bullets[idx : idx + size])
            idx += size
        return chunks
    # 불렛 수 허용 범위(n ≤ max): 줄 수 기반 판단 (+1 여유)
    if total_lines <= _MAX_LINES_SLIDE + 1:
        return [bullets]
    # greedy 줄 수 기반 분할
    target = _MAX_LINES_SLIDE - 2
    chunks, current, curr_lines = [], [], 0
    for b in bullets:
        bl = _est_lines(b)
        if current and curr_lines + bl > target:
            chunks.append(current)
            current, curr_lines = [b], bl
        else:
            current.append(b)
            curr_lines += bl
    if current:
        chunks.append(current)
    return chunks if len(chunks) > 1 else [bullets]


def build_after(data: dict, output_path: str, cover_image: str = None):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    t = data["title"]
    build_title_slide(prs, t["series"], t["lecture_title"], t["pastor"],
                      cover_image=cover_image)

    for cs in data["content_slides"]:
        heading = cs["heading"]
        bullets = cs["bullets"]
        images  = cs.get("images", [])

        if not bullets:
            build_content_slide(prs, heading, [], images=images)
            continue

        for j, chunk in enumerate(_split_bullets(bullets, MAX_BULLETS)):
            # 이미지는 첫 번째 청크에만 포함
            build_content_slide(prs, heading, chunk,
                                images=(images if j == 0 else None),
                                is_continued=(j > 0))

    prs.save(output_path)
    return len(prs.slides)


# ── PNG 내보내기 (PowerPoint COM, Windows 전용) ────────────────────────────────

def export_png(pptx_path: str, out_dir: str, width_px=1920, height_px=1080):
    import comtypes.client
    pptx_path = os.path.abspath(pptx_path)
    out_dir   = os.path.abspath(out_dir)
    os.makedirs(out_dir, exist_ok=True)

    app = comtypes.client.CreateObject("PowerPoint.Application")
    app.Visible = 1
    try:
        prs = app.Presentations.Open(pptx_path, ReadOnly=True, WithWindow=False)
        try:
            for i, slide in enumerate(prs.Slides, start=1):
                png_path = os.path.join(out_dir, f"slide_{i:03d}.png")
                slide.Export(png_path, "PNG", width_px, height_px)
                print(f"  PNG: {os.path.basename(png_path)}")
        finally:
            prs.Close()
    finally:
        app.Quit()

    print(f"  완료 → {out_dir}")


# ── 메인 ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="PL Church PPT 변환기")
    parser.add_argument("input",           help="Before PPTX 경로")
    parser.add_argument("output", nargs="?", help="출력 경로 (생략 시 자동)")
    parser.add_argument("--png", action="store_true",
                        help="변환 후 PNG 내보내기 (PowerPoint 필요)")
    parser.add_argument("--png-dir",       help="PNG 저장 폴더")
    parser.add_argument("--canva", action="store_true",
                        help="변환 후 Canva에 자동 업로드 (canva_api.py 인증 필요)")
    args = parser.parse_args()

    if not os.path.exists(args.input):
        print(f"오류: 파일 없음 → {args.input}")
        sys.exit(1)

    out_path = args.output or (os.path.splitext(args.input)[0] + "_converted.pptx")

    print(f"[1/3] 파싱: {os.path.basename(args.input)}")
    data = parse_before(args.input)
    t = data["title"]
    print(f"      시리즈: {t['series']}")
    print(f"      제목  : {t['lecture_title']}")
    print(f"      목사  : {t['pastor']}")
    print(f"      내용 슬라이드: {len(data['content_slides'])}개 (분리 전)")

    # 표지 이미지: cover.png → 표지.png 순으로 탐색
    script_dir  = os.path.dirname(os.path.abspath(__file__))
    cover_image = next(
        (os.path.join(script_dir, n) for n in ("cover.png", "표지.png")
         if os.path.exists(os.path.join(script_dir, n))),
        None,
    )
    if cover_image:
        print(f"[2/3] 변환 중... (표지 이미지 적용)")
    else:
        print(f"[2/3] 변환 중... (표지 이미지 없음, 기본 배경 사용)")

    n = build_after(data, out_path, cover_image=cover_image)
    print(f"[3/3] 완료 → {os.path.basename(out_path)}  ({n} 슬라이드)")

    if args.png:
        png_dir = args.png_dir or os.path.splitext(out_path)[0] + "_slides"
        print(f"[+]  PNG 내보내기 → {png_dir}")
        export_png(out_path, png_dir)

    if args.canva:
        try:
            from canva_api import upload_pptx
            upload_pptx(out_path, verbose=True)
        except ImportError:
            print("오류: canva_api.py 파일이 없습니다.")
        except RuntimeError as e:
            print(f"Canva 업로드 오류: {e}")


if __name__ == "__main__":
    main()
