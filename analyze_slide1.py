import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation

def analyze_slide1(path):
    prs = Presentation(path)
    slide = list(prs.slides)[0]
    print(f"\n=== {path} ===")
    for j, shape in enumerate(slide.shapes):
        print(f"  Shape[{j}] name={shape.name!r}")
        if shape.has_text_frame:
            for k, para in enumerate(shape.text_frame.paragraphs):
                t = "".join(r.text for r in para.runs)
                if t.strip():
                    print(f"    para[{k}]: {t!r}")

base = r"C:\Users\myc43\OneDrive - ETERNAL LIBRTY POLICY INSTITUTE\바탕 화면\plc ppt"
for f in ["에베소서 제4강(before).pptx", "레위기 제12강(before).pptx", "마태복음 제41강(before).pptx"]:
    analyze_slide1(f"{base}\\{f}")
