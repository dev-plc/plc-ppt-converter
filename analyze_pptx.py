#!/usr/bin/env python3
"""
Detailed PPTX analysis script.
Analyzes Original, Converted, and Complete PPTX files.
"""
import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from lxml import etree

BASE = "/root/.claude/uploads/13f911f1-87ab-4d0f-85c1-4235e95894e6"
ORIG_PATH  = f"{BASE}/432e161a-_____1_Original.pptx"
CONV_PATH  = f"{BASE}/ce5cf320-_____1__converted_7.pptx"
COMP_PATH  = f"{BASE}/4adf2458-_____1_Complete.pptx"

NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

def rgb_hex(color_obj):
    """Return hex string or None."""
    try:
        if color_obj and color_obj.type is not None:
            rgb = color_obj.rgb
            return f"#{rgb.red:02X}{rgb.green:02X}{rgb.blue:02X}"
    except Exception:
        pass
    return None

def run_detail(run):
    font = run.font
    color = rgb_hex(font.color)
    size = int(font.size.pt) if font.size else None
    # Check highlight in XML
    hl_val = None
    try:
        rpr = run._r.find(f"{{{NS}}}rPr")
        if rpr is not None:
            hl = rpr.find(f"{{{NS}}}highlight")
            if hl is not None:
                hl_val = hl.get("val")
    except Exception:
        pass
    return {
        "text": run.text,
        "bold": font.bold,
        "italic": font.italic,
        "underline": font.underline,
        "color": color,
        "size": size,
        "highlight": hl_val,
    }

def format_run(r):
    parts = []
    if r["bold"]: parts.append("BOLD")
    if r["underline"]: parts.append("UL")
    if r["color"]: parts.append(f"color={r['color']}")
    if r["highlight"]: parts.append(f"HL={r['highlight']}")
    if r["size"]: parts.append(f"{r['size']}pt")
    desc = ", ".join(parts)
    return f"[{r['text']!r} | {desc}]"

# ─────────────────────────────────────────────────────────────
# SECTION 1: Slide counts
# ─────────────────────────────────────────────────────────────
print("=" * 70)
print("SECTION 1: SLIDE COUNTS")
print("=" * 70)

prs_orig = Presentation(ORIG_PATH)
prs_conv = Presentation(CONV_PATH)
prs_comp = Presentation(COMP_PATH)

print(f"Original  : {len(prs_orig.slides)} slides")
print(f"Converted : {len(prs_conv.slides)} slides")
print(f"Complete  : {len(prs_comp.slides)} slides")

# ─────────────────────────────────────────────────────────────
# SECTION 2: Original slide 9 — all shapes, full run detail
# ─────────────────────────────────────────────────────────────
print()
print("=" * 70)
print("SECTION 2: ORIGINAL SLIDE 9 — FULL SHAPE & RUN ANALYSIS")
print("=" * 70)

slide9 = prs_orig.slides[8]  # 0-indexed
for si, shape in enumerate(slide9.shapes):
    print(f"\n  [Shape {si}] name='{shape.name}'  has_tf={shape.has_text_frame}")
    if shape.has_text_frame:
        for pi, para in enumerate(shape.text_frame.paragraphs):
            para_text = para.text
            if not para_text.strip() and not para.runs:
                print(f"    Para {pi}: <empty>")
                continue
            print(f"    Para {pi}: '{para_text}'")
            for ri, run in enumerate(para.runs):
                info = run_detail(run)
                print(f"      Run {ri}: {format_run(info)}")

# ─────────────────────────────────────────────────────────────
# SECTION 3: Converted vs Complete — ALL slides side by side
# ─────────────────────────────────────────────────────────────
print()
print("=" * 70)
print("SECTION 3: CONVERTED vs COMPLETE — SLIDE-BY-SLIDE")
print("=" * 70)

def slide_summary(slide, label, idx):
    print(f"  [{label} Slide {idx+1}]")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        for para in tf.paragraphs:
            t = para.text.strip()
            if not t:
                continue
            runs_str = "  ".join(format_run(run_detail(r)) for r in para.runs if r.text)
            print(f"    [{shape.name}] '{t}'")
            if runs_str:
                print(f"      {runs_str}")

n_conv = len(prs_conv.slides)
n_comp = len(prs_comp.slides)
max_n = max(n_conv, n_comp)

for idx in range(max_n):
    print(f"\n{'─'*70}")
    print(f"  SLIDE PAIR {idx+1}")
    print(f"{'─'*70}")
    if idx < n_conv:
        slide_summary(prs_conv.slides[idx], "CONVERTED", idx)
    else:
        print(f"  [CONVERTED Slide {idx+1}]: *** DOES NOT EXIST ***")
    print()
    if idx < n_comp:
        slide_summary(prs_comp.slides[idx], "COMPLETE", idx)
    else:
        print(f"  [COMPLETE  Slide {idx+1}]: *** DOES NOT EXIST ***")

# ─────────────────────────────────────────────────────────────
# SECTION 4: Complete — color pattern summary
# ─────────────────────────────────────────────────────────────
print()
print("=" * 70)
print("SECTION 4: COMPLETE VERSION — GOLD vs WHITE COLOR PATTERN")
print("=" * 70)

for si, slide in enumerate(prs_comp.slides):
    gold_texts = []
    white_texts = []
    other_texts = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                t = run.text.strip()
                if not t:
                    continue
                c = rgb_hex(run.font.color)
                if c and c.upper() == "#F5C518":
                    gold_texts.append(t)
                elif c and c.upper() in ("#FFFFFF", "#FEFEFE"):
                    white_texts.append(t)
                elif c:
                    other_texts.append((t, c))
    print(f"\n  Slide {si+1}:")
    if gold_texts:
        print(f"    GOLD : {gold_texts}")
    if white_texts:
        print(f"    WHITE: {white_texts}")
    if other_texts:
        print(f"    OTHER: {other_texts}")
    if not gold_texts and not white_texts and not other_texts:
        print(f"    (no colored runs)")

# ─────────────────────────────────────────────────────────────
# SECTION 5: Complete — bullet count & split pattern
# ─────────────────────────────────────────────────────────────
print()
print("=" * 70)
print("SECTION 5: COMPLETE VERSION — BULLET COUNT PER SLIDE")
print("=" * 70)

for si, slide in enumerate(prs_comp.slides):
    all_paras = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t:
                all_paras.append(f"[{shape.name}]'{t}'")
    print(f"  Slide {si+1} ({len(all_paras)} paras): {' | '.join(all_paras)}")

print("\nDone.")
